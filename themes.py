from pptx import Presentation
from pptx.util import Inches, Pt
# importing modules
from os import startfile
import io
import requests
import urllib.request 
from bs4 import BeautifulSoup
import re
from requests_html import HTMLSession
session = HTMLSession()

def create_ppt_1(topic, author_main, detail_slide, theme="theme_0"):
	def google_images(query):
	    links = []
	    response = session.get(f'https://www.google.com/search?q={query}&tbm=isch&sclient=img')
	    res = (response.text)
	    a = re.findall(r'(https?://[^\s]+)', res)
	    for item in a:
	        if (".jpg") in item :
	            b = (item.partition(".jpg")[0])
	            c = b.partition(""",["https""")[2]
	            d = "https"+c+".jpg"
	            for item in c:
	                if c != "":
	                    if d not in links:
	                        if ".jpg" in d:
	                            links.append(d)
	    return(links)
	def remove_brackets(string):
		string=re.sub("\[.*?\]","[]",string)
		string=re.sub("\(.*?\)","",string)
		string = string.replace("[]", "")
		return string		
	def get_data(query):
	    query = query.replace(" ", "_")
	    data_list =[("h3", "Introduction")];image_list =[]
	    html = urllib.request.urlopen("https://en.wikipedia.org/wiki/" +  query)
	    htmlParse = BeautifulSoup(html, 'html.parser')
	    intro_found = False
	    for para in htmlParse.find_all(["img", "p", "h3"]):
	        a, b = 0, 0
	        try:
	            if not ".svg" in para['src'] and not "footer" in para['src']:
	                if ".png" in para['src'] or ".jpg" in para['src']:
	                    a = (para.name, para['src'])
	        except Exception:
	        	b = (para.name, remove_brackets(para.get_text()))
	        	if para.name =="p":
	        		b = (para.name, remove_brackets(para.get_text())[:-2])
	        if a:
	            image_list.append(a)
	        if b:
	        	if intro_found:
	        		data_list.append(b)
	        	else:
		        	if len(str(b[1]).split(" "))>5:
			            data_list.append(b)
			            intro_found = True
	    i = 0
	    main_list=[]
	    while not i>= len(data_list)-1:
	        if data_list[i][0] == 'h3' and data_list[i+1][0] == 'p':
	            main_list.append((data_list[i], data_list[i+1]))
	            i+=2
	        else:
	            i+=1
	    return(main_list)

	data = get_data(topic.title())
	image_list = google_images(topic)
	if len(data)<detail_slide:
		detail_slide=len(data)
	if detail_slide>10:
		detail_slide=10
	title,author = topic.title(), "Created by "+author_main
	index_topics = "1. Start\n2. Middle\n3. End of the Index"
	#theme1
	ppt = Presentation(pptx=f"{theme}.pptx")

	#SLIDE 1


	def slide(slide_num, slide_title, paragraph):
		slide_1 = ppt.slides[slide_num]
		slide_1.shapes[0].text_frame.paragraphs[0].runs[0].text = slide_title
		slide_1.shapes[1].text_frame.paragraphs[0].runs[0].text = paragraph
	slide(0,title, author)
	def generate_index():
		index = "Acknowledgement\n"
		for i in range(detail_slide):
			index+= data[i][0][1]+"\n"
		return index+"Bilbiography\nThankyou"
	slide(1,"INDEX", generate_index())
	global image_skip
	image_skip = 0
	image_upload = True
	def change_image(slide, img_index):
		global image_skip
		image_upload = True
		while image_upload:
			response = requests.get((image_list[img_index+image_skip]))
			f = bytearray(response.content)
			if str(response.status_code).startswith('4') or len(f)<1500:
				image_skip+=1
			else: 
				shape = slide.shapes[2]
				slide_part, rId = shape.part, shape._element.blip_rId
				image_part = slide_part.related_part(rId)
				image_part.blob = f
				image_upload = False

	change_image(ppt.slides[0], 0)
	def generate_topic():
		for i in range(3,3+detail_slide):
			slide= ppt.slides[i]
			change_image(slide, i-2)
			slide.shapes[0].text_frame.paragraphs[0].runs[0].text = str(i-2)+". "+data[i-3][0][1]
			slide.shapes[1].text_frame.paragraphs[0].runs[0].text = data[i-3][1][1]
			if len(data[i-3][1][1])>1000:
				slide.shapes[1].text_frame.paragraphs[0].runs[0].font.size = Pt(16)
			elif len(data[i-3][1][1])>600:
				slide.shapes[1].text_frame.paragraphs[0].runs[0].font.size = Pt(18)


	generate_topic()

	def delete_slide(presentation,  index):
	    xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
	    slides = list(xml_slides)
	    xml_slides.remove(slides[index])
	for i in range(3+detail_slide, 13):
		delete_slide(ppt, 3+detail_slide)

	ppt.save("sample.pptx")
	startfile("sample.pptx")


def create_ppt_2(topic, author_main, detail_slide, theme="theme_1"):
	def google_images(query):
	    links = []
	    response = session.get(f'https://www.google.com/search?q={query}&tbm=isch&sclient=img')
	    res = (response.text)
	    a = re.findall(r'(https?://[^\s]+)', res)
	    for item in a:
	        if (".jpg") in item :
	            b = (item.partition(".jpg")[0])
	            c = b.partition(""",["https""")[2]
	            d = "https"+c+".jpg"
	            for item in c:
	                if c != "":
	                    if d not in links:
	                        if ".jpg" in d:
	                            links.append(d)
	    return(links)
	def remove_brackets(string):
		string=re.sub("\[.*?\]","[]",string)
		string=re.sub("\(.*?\)","",string)
		string = string.replace("[]", "")
		return string		
	def get_data(query):
	    query = query.replace(" ", "_")
	    data_list =[("h3", "Introduction")];image_list =[]
	    html = urllib.request.urlopen("https://en.wikipedia.org/wiki/" +  query)
	    htmlParse = BeautifulSoup(html, 'html.parser')
	    intro_found = False
	    for para in htmlParse.find_all(["img", "p", "h3"]):
	        a, b = 0, 0
	        try:
	            if not ".svg" in para['src'] and not "footer" in para['src']:
	                if ".png" in para['src'] or ".jpg" in para['src']:
	                    a = (para.name, para['src'])
	        except Exception:
	            b = (para.name, remove_brackets(para.get_text()))
	            if para.name =="p":
	            	b = (para.name, remove_brackets(para.get_text())[:-2])
	        if a:
	            image_list.append(a)
	        if b:
	        	if intro_found:
	        		data_list.append(b)
	        	else:
		        	if len(str(b[1]).split(" "))>5:
			            data_list.append(b)
			            intro_found = True
	    i = 0
	    main_list=[]
	    while not i>= len(data_list)-1:
	        if data_list[i][0] == 'h3' and data_list[i+1][0] == 'p':
	            main_list.append((data_list[i], data_list[i+1]))
	            i+=2
	        else:
	            i+=1
	    return(main_list)

	data = get_data(topic.title())
	image_list = google_images(topic)
	if len(data)<detail_slide:
		detail_slide=len(data)
	if detail_slide>10:
		detail_slide=10
	title,author = topic.title(), "Created by "+author_main
	index_topics = "1. Start\n2. Middle\n3. End of the Index"
	#theme1
	ppt = Presentation(pptx=f"{theme}.pptx")

	#SLIDE 1


	def slide(slide_num, slide_title, paragraph, slide_title_pos, para_pos):
		slide_1 = ppt.slides[slide_num]
		slide_1.shapes[slide_title_pos].text_frame.paragraphs[0].runs[0].text = slide_title
		slide_1.shapes[para_pos].text_frame.paragraphs[0].runs[0].text = paragraph
	slide(0,title, author, 2, 3)
	def generate_index():
		index = "Acknowledgement\n"
		for i in range(detail_slide):
			index+= data[i][0][1]+"\n"
		return index+"Bilbiography\nThankyou"
	slide(1,"INDEX", generate_index(), 4, 5)
	global image_skip
	image_skip = 0
	image_upload = True
	def change_image(slide, img_index, pos=4):
		global image_skip
		image_upload = True
		while image_upload:
			response = requests.get((image_list[img_index+image_skip]))
			f = bytearray(response.content)
			if str(response.status_code).startswith('4') or len(f)<1500:
				image_skip+=1
			else: 
				shape = slide.shapes[pos]
				slide_part, rId = shape.part, shape._element.blip_rId
				image_part = slide_part.related_part(rId)
				image_part.blob = f
				image_upload = False

	change_image(ppt.slides[0], 0, 5)
	def generate_topic():
		for i in range(3,3+detail_slide):
			slide= ppt.slides[i]
			change_image(slide, i-2)
			slide.shapes[1].text_frame.paragraphs[0].runs[0].text = str(i-2)+". "+data[i-3][0][1]
			slide.shapes[3].text_frame.paragraphs[0].runs[0].text = data[i-3][1][1]
			if len(data[i-3][1][1])>1000:
				slide.shapes[3].text_frame.paragraphs[0].runs[0].font.size = Pt(16)
			elif len(data[i-3][1][1])>600:
				slide.shapes[3].text_frame.paragraphs[0].runs[0].font.size = Pt(18)


	generate_topic()

	def delete_slide(presentation,  index):
	    xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
	    slides = list(xml_slides)
	    xml_slides.remove(slides[index])
	for i in range(3+detail_slide, 13):
		delete_slide(ppt, 3+detail_slide)

	ppt.save("sample.pptx")
	startfile("sample.pptx")


def create_ppt_3(topic, author_main, detail_slide, theme="theme_2"):
	def google_images(query):
	    links = []
	    response = session.get(f'https://www.google.com/search?q={query}&tbm=isch&sclient=img')
	    res = (response.text)
	    a = re.findall(r'(https?://[^\s]+)', res)
	    for item in a:
	        if (".jpg") in item :
	            b = (item.partition(".jpg")[0])
	            c = b.partition(""",["https""")[2]
	            d = "https"+c+".jpg"
	            for item in c:
	                if c != "":
	                    if d not in links:
	                        if ".jpg" in d:
	                            links.append(d)
	    return(links)
	def remove_brackets(string):
		string=re.sub("\[.*?\]","[]",string)
		string=re.sub("\(.*?\)","",string)
		string = string.replace("[]", "")
		return string		
	def get_data(query):
	    query = query.replace(" ", "_")
	    data_list =[("h3", "Introduction")];image_list =[]
	    html = urllib.request.urlopen("https://en.wikipedia.org/wiki/" +  query)
	    htmlParse = BeautifulSoup(html, 'html.parser')
	    intro_found = False
	    for para in htmlParse.find_all(["img", "p", "h3"]):
	        a, b = 0, 0
	        try:
	            if not ".svg" in para['src'] and not "footer" in para['src']:
	                if ".png" in para['src'] or ".jpg" in para['src']:
	                    a = (para.name, para['src'])
	        except Exception:
	        	b = (para.name, remove_brackets(para.get_text()))
	        	if para.name =="p":
	        		b = (para.name, remove_brackets(para.get_text())[:-2])
	        if a:
	            image_list.append(a)
	        if b:
	        	if intro_found:
	        		data_list.append(b)
	        	else:
		        	if len(str(b[1]).split(" "))>5:
			            data_list.append(b)
			            intro_found = True
	    i = 0
	    main_list=[]
	    while not i>= len(data_list)-1:
	        if data_list[i][0] == 'h3' and data_list[i+1][0] == 'p':
	            main_list.append((data_list[i], data_list[i+1]))
	            i+=2
	        else:
	            i+=1
	    return(main_list)

	data = get_data(topic.title())
	image_list = google_images(topic)
	if len(data)<detail_slide:
		detail_slide=len(data)
	if detail_slide>10:
		detail_slide=10
	title,author = topic.title(), "Created by "+author_main
	index_topics = "1. Start\n2. Middle\n3. End of the Index"
	#theme1
	ppt = Presentation(pptx=f"{theme}.pptx")

	#SLIDE 1


	def slide(slide_num, slide_title, paragraph):
		slide_1 = ppt.slides[slide_num]
		slide_1.shapes[0].text_frame.paragraphs[0].runs[0].text = slide_title
		slide_1.shapes[1].text_frame.paragraphs[0].runs[0].text = paragraph
	slide(0,title, author)
	def generate_index():
		index = "Acknowledgement\n"
		for i in range(detail_slide):
			index+= data[i][0][1]+"\n"
		return index+"Bilbiography\nThankyou"
	slide(1,"INDEX", generate_index())
	global image_skip
	image_skip = 0
	image_upload = True
	def change_image(slide, img_index):
		global image_skip
		image_upload = True
		while image_upload:
			response = requests.get((image_list[img_index+image_skip]))
			f = bytearray(response.content)
			if str(response.status_code).startswith('4') or len(f)<1500:
				image_skip+=1
			else: 
				shape = slide.shapes[2]
				slide_part, rId = shape.part, shape._element.blip_rId
				image_part = slide_part.related_part(rId)
				image_part.blob = f
				image_upload = False

	change_image(ppt.slides[0], 0)
	def generate_topic():
		for i in range(3,3+detail_slide):
			slide= ppt.slides[i]
			change_image(slide, i-2)
			slide.shapes[0].text_frame.paragraphs[0].runs[0].text = str(i-2)+". "+data[i-3][0][1]
			slide.shapes[1].text_frame.paragraphs[0].runs[0].text = data[i-3][1][1]
			if len(data[i-3][1][1])>1000:
				slide.shapes[1].text_frame.paragraphs[0].runs[0].font.size = Pt(16)
			elif len(data[i-3][1][1])>600:
				slide.shapes[1].text_frame.paragraphs[0].runs[0].font.size = Pt(18)


	generate_topic()

	def delete_slide(presentation,  index):
	    xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
	    slides = list(xml_slides)
	    xml_slides.remove(slides[index])
	for i in range(3+detail_slide, 13):
		delete_slide(ppt, 3+detail_slide)

	ppt.save("sample.pptx")
	startfile("sample.pptx")



def create_ppt_4(topic, author_main, detail_slide, theme="theme_3"):
	def google_images(query):
	    links = []
	    response = session.get(f'https://www.google.com/search?q={query}&tbm=isch&sclient=img')
	    res = (response.text)
	    a = re.findall(r'(https?://[^\s]+)', res)
	    for item in a:
	        if (".jpg") in item :
	            b = (item.partition(".jpg")[0])
	            c = b.partition(""",["https""")[2]
	            d = "https"+c+".jpg"
	            for item in c:
	                if c != "":
	                    if d not in links:
	                        if ".jpg" in d:
	                            links.append(d)
	    return(links)
	def remove_brackets(string):
		string=re.sub("\[.*?\]","[]",string)
		string=re.sub("\(.*?\)","",string)
		string = string.replace("[]", "")
		return string		
	def get_data(query):
	    query = query.replace(" ", "_")
	    data_list =[("h3", "Introduction")];image_list =[]
	    html = urllib.request.urlopen("https://en.wikipedia.org/wiki/" +  query)
	    htmlParse = BeautifulSoup(html, 'html.parser')
	    intro_found = False
	    for para in htmlParse.find_all(["img", "p", "h3"]):
	        a, b = 0, 0
	        try:
	            if not ".svg" in para['src'] and not "footer" in para['src']:
	                if ".png" in para['src'] or ".jpg" in para['src']:
	                    a = (para.name, para['src'])
	        except Exception:
	        	b = (para.name, remove_brackets(para.get_text()))
	        	if para.name =="p":
	        		b = (para.name, remove_brackets(para.get_text())[:-2])
	        if a:
	            image_list.append(a)
	        if b:
	        	if intro_found:
	        		data_list.append(b)
	        	else:
		        	if len(str(b[1]).split(" "))>5:
			            data_list.append(b)
			            intro_found = True
	    i = 0
	    main_list=[]
	    while not i>= len(data_list)-1:
	        if data_list[i][0] == 'h3' and data_list[i+1][0] == 'p':
	            main_list.append((data_list[i], data_list[i+1]))
	            i+=2
	        else:
	            i+=1
	    return(main_list)

	data = get_data(topic.title())
	image_list = google_images(topic)
	if len(data)<detail_slide:
		detail_slide=len(data)
	if detail_slide>10:
		detail_slide=10
	title,author = topic.title(), "Created by "+author_main
	index_topics = "1. Start\n2. Middle\n3. End of the Index"
	#theme1
	ppt = Presentation(pptx=f"{theme}.pptx")

	#SLIDE 1


	def slide(slide_num, slide_title, paragraph):
		slide_1 = ppt.slides[slide_num]
		slide_1.shapes[0].text_frame.paragraphs[0].runs[0].text = slide_title
		slide_1.shapes[1].text_frame.paragraphs[0].runs[0].text = paragraph
	slide(0,title, author)
	def generate_index():
		index = "Acknowledgement\n"
		for i in range(detail_slide):
			index+= data[i][0][1]+"\n"
		return index+"Bilbiography\nThankyou"
	slide(1,"INDEX", generate_index())
	global image_skip
	image_skip = 0
	image_upload = True
	def change_image(slide, img_index):
		global image_skip
		image_upload = True
		while image_upload:
			response = requests.get((image_list[img_index+image_skip]))
			f = bytearray(response.content)
			if str(response.status_code).startswith('4') or len(f)<1500:
				image_skip+=1
			else: 
				shape = slide.shapes[2]
				slide_part, rId = shape.part, shape._element.blip_rId
				image_part = slide_part.related_part(rId)
				image_part.blob = f
				image_upload = False

	change_image(ppt.slides[0], 0)
	def generate_topic():
		for i in range(3,3+detail_slide):
			slide= ppt.slides[i]
			change_image(slide, i-2)
			slide.shapes[0].text_frame.paragraphs[0].runs[0].text = str(i-2)+". "+data[i-3][0][1]
			slide.shapes[1].text_frame.paragraphs[0].runs[0].text = data[i-3][1][1]
			if len(data[i-3][1][1])>1000:
				slide.shapes[1].text_frame.paragraphs[0].runs[0].font.size = Pt(16)
			elif len(data[i-3][1][1])>600:
				slide.shapes[1].text_frame.paragraphs[0].runs[0].font.size = Pt(18)


	generate_topic()

	def delete_slide(presentation,  index):
	    xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
	    slides = list(xml_slides)
	    xml_slides.remove(slides[index])
	for i in range(3+detail_slide, 13):
		delete_slide(ppt, 3+detail_slide)

	ppt.save("sample.pptx")
	startfile("sample.pptx")



def create_ppt_5(topic, author_main, detail_slide, theme="theme_4"):
	def google_images(query):
	    links = []
	    response = session.get(f'https://www.google.com/search?q={query}&tbm=isch&sclient=img')
	    res = (response.text)
	    a = re.findall(r'(https?://[^\s]+)', res)
	    for item in a:
	        if (".jpg") in item :
	            b = (item.partition(".jpg")[0])
	            c = b.partition(""",["https""")[2]
	            d = "https"+c+".jpg"
	            for item in c:
	                if c != "":
	                    if d not in links:
	                        if ".jpg" in d:
	                            links.append(d)
	    return(links)
	def remove_brackets(string):
		string=re.sub("\[.*?\]","[]",string)
		string=re.sub("\(.*?\)","",string)
		string = string.replace("[]", "")
		return string		
	def get_data(query):
	    query = query.replace(" ", "_")
	    data_list =[("h3", "Introduction")];image_list =[]
	    html = urllib.request.urlopen("https://en.wikipedia.org/wiki/" +  query)
	    htmlParse = BeautifulSoup(html, 'html.parser')
	    intro_found = False
	    for para in htmlParse.find_all(["img", "p", "h3"]):
	        a, b = 0, 0
	        try:
	            if not ".svg" in para['src'] and not "footer" in para['src']:
	                if ".png" in para['src'] or ".jpg" in para['src']:
	                    a = (para.name, para['src'])
	        except Exception:
	            b = (para.name, remove_brackets(para.get_text()))
	            if para.name =="p":
	            	b = (para.name, remove_brackets(para.get_text())[:-2])
	        if a:
	            image_list.append(a)
	        if b:
	        	if intro_found:
	        		data_list.append(b)
	        	else:
		        	if len(str(b[1]).split(" "))>5:
			            data_list.append(b)
			            intro_found = True
	    i = 0
	    main_list=[]
	    while not i>= len(data_list)-1:
	        if data_list[i][0] == 'h3' and data_list[i+1][0] == 'p':
	            main_list.append((data_list[i], data_list[i+1]))
	            i+=2
	        else:
	            i+=1
	    return(main_list)

	data = get_data(topic.title())
	image_list = google_images(topic)
	if len(data)<detail_slide:
		detail_slide=len(data)
	if detail_slide>10:
		detail_slide=10
	title,author = topic.title(), "Created by "+author_main
	index_topics = "1. Start\n2. Middle\n3. End of the Index"
	#theme1
	ppt = Presentation(pptx=f"{theme}.pptx")

	#SLIDE 1


	def slide(slide_num, slide_title, paragraph, slide_title_pos, para_pos):
		slide_1 = ppt.slides[slide_num]
		slide_1.shapes[slide_title_pos].text_frame.paragraphs[0].runs[0].text = slide_title
		slide_1.shapes[para_pos].text_frame.paragraphs[0].runs[0].text = paragraph
	slide(0,title, author, 0, 1)
	def generate_index():
		index = "Acknowledgement\n"
		for i in range(detail_slide):
			index+= data[i][0][1]+"\n"
		return index+"Bilbiography\nThankyou"
	slide(1,"INDEX", generate_index(), 5, 6)
	global image_skip
	image_skip = 0
	image_upload = True
	def change_image(slide, img_index, pos=3):
		global image_skip
		image_upload = True
		while image_upload:
			response = requests.get((image_list[img_index+image_skip]))
			f = bytearray(response.content)
			if str(response.status_code).startswith('4') or len(f)<1500:
				image_skip+=1
			else: 
				shape = slide.shapes[pos]
				slide_part, rId = shape.part, shape._element.blip_rId
				image_part = slide_part.related_part(rId)
				image_part.blob = f
				image_upload = False

	#change_image(ppt.slides[0], 0, 2)
	def generate_topic():
		for i in range(3,3+detail_slide):
			slide= ppt.slides[i]
			change_image(slide, i-2)
			slide.shapes[4].text_frame.paragraphs[0].runs[0].text = str(i-2)+". "+data[i-3][0][1]
			slide.shapes[6].text_frame.paragraphs[0].runs[0].text = data[i-3][1][1]
			if len(data[i-3][1][1])>1000:
				slide.shapes[6].text_frame.paragraphs[0].runs[0].font.size = Pt(16)
			elif len(data[i-3][1][1])>600:
				slide.shapes[6].text_frame.paragraphs[0].runs[0].font.size = Pt(18)


	generate_topic()

	def delete_slide(presentation,  index):
	    xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
	    slides = list(xml_slides)
	    xml_slides.remove(slides[index])
	for i in range(3+detail_slide, 13):
		delete_slide(ppt, 3+detail_slide)

	ppt.save("sample.pptx")
	startfile("sample.pptx")

#create_ppt_1("niall horan", 'hamza', 4)